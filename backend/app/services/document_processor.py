import uuid
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple

from app.config import settings

logger = logging.getLogger(__name__)

# Optional heavy dependencies
try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    LANGCHAIN_AVAILABLE = True
except ImportError:
    LANGCHAIN_AVAILABLE = False

class DocumentProcessor:
    """Handles document upload, parsing, and chunking"""

    def __init__(self):
        if LANGCHAIN_AVAILABLE:
            self.text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=settings.CHUNK_SIZE,
                chunk_overlap=settings.CHUNK_OVERLAP,
                length_function=len,
                separators=["\n\n", "\n", ". ", " ", ""]
            )
        else:
            self.text_splitter = None

    async def process_document(
        self,
        file_path: Path,
        filename: str,
        document_id: Optional[str] = None,
    ) -> Dict[str, Any]:
        """Process uploaded document and return chunks with metadata"""

        file_extension = file_path.suffix.lower()
        document_id = document_id or str(uuid.uuid4())

        # Extract text based on file type
        if file_extension == ".pdf":
            if not PYPDF2_AVAILABLE:
                raise ValueError("PDF processing requires PyPDF2. Install it with: pip install PyPDF2")
            text, metadata = self._extract_pdf(file_path)
        elif file_extension in [".pptx", ".ppt"]:
            if not PPTX_AVAILABLE:
                raise ValueError("PPTX processing requires python-pptx. Install it with: pip install python-pptx")
            text, metadata = self._extract_pptx(file_path)
        elif file_extension in [".docx", ".doc"]:
            if not DOCX_AVAILABLE:
                raise ValueError("DOCX processing requires python-docx. Install it with: pip install python-docx")
            text, metadata = self._extract_docx(file_path)
        elif file_extension == ".txt":
            text, metadata = self._extract_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")

        # Chunk the text
        if self.text_splitter:
            chunks = self.text_splitter.split_text(text)
        else:
            # Simple fallback chunking
            chunk_size = settings.CHUNK_SIZE
            chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size - settings.CHUNK_OVERLAP)]

        # Prepare chunk metadata
        chunk_data = []
        for idx, chunk in enumerate(chunks):
            chunk_data.append({
                "text": chunk,
                "metadata": {
                    "document_id": document_id,
                    "filename": filename,
                    "chunk_index": idx,
                    "total_chunks": len(chunks),
                    **metadata
                }
            })

        return {
            "document_id": document_id,
            "filename": filename,
            "chunks": chunk_data,
            "metadata": metadata
        }

    def _extract_pdf(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PDF"""
        text = ""
        page_count = 0

        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            page_count = len(pdf_reader.pages)

            for page_num, page in enumerate(pdf_reader.pages, start=1):
                page_text = page.extract_text()
                text += f"\n--- Page {page_num} ---\n{page_text}\n"

        metadata = {
            "file_type": "pdf",
            "total_pages": page_count
        }

        return text, metadata

    def _extract_pptx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PowerPoint"""
        prs = Presentation(file_path)
        text = ""
        slide_count = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = f"\n--- Slide {slide_num} ---\n"

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"

            text += slide_text

        metadata = {
            "file_type": "pptx",
            "total_pages": slide_count
        }

        return text, metadata

    def _extract_docx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from Word document"""
        doc = DocxDocument(file_path)
        text = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])

        metadata = {
            "file_type": "docx",
            "total_pages": None  # Word doesn't have fixed pages
        }

        return text, metadata

    def _extract_txt(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from plain text file"""
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()

        metadata = {
            "file_type": "txt",
            "total_pages": None
        }

        return text, metadata
